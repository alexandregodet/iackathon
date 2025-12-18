#!/usr/bin/env python3
"""
Script to create a Kaamelott-themed PowerPoint presentation
for the IAckathon project features.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Kaamelott color palette
PARCHMENT = RgbColor(0xF4, 0xE4, 0xBC)
PARCHMENT_DARK = RgbColor(0xE8, 0xD4, 0xA8)
BURGUNDY = RgbColor(0x72, 0x2F, 0x37)
BURGUNDY_DARK = RgbColor(0x5A, 0x25, 0x2C)
GOLD = RgbColor(0xC9, 0xA2, 0x27)
FOREST = RgbColor(0x2D, 0x4A, 0x3E)
INK = RgbColor(0x2C, 0x18, 0x10)
INK_LIGHT = RgbColor(0x4A, 0x37, 0x28)

def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=None):
    """Add a title slide with Kaamelott styling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, PARCHMENT)

    # Add burgundy banner at top
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(2)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = BURGUNDY
    banner.line.fill.background()

    # Gold line under banner
    gold_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.95),
        Inches(13.333), Inches(0.1)
    )
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = GOLD
    gold_line.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(12.333), Inches(1.3)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3),
            Inches(11.333), Inches(1.5)
        )
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = INK
        p.font.name = "Georgia"
        p.alignment = PP_ALIGN.CENTER

    # Decorative fleur-de-lis
    deco_box = slide.shapes.add_textbox(
        Inches(6), Inches(5),
        Inches(1.333), Inches(0.8)
    )
    tf = deco_box.text_frame
    p = tf.paragraphs[0]
    p.text = "\u269C"  # Fleur-de-lis
    p.font.size = Pt(48)
    p.font.color.rgb = BURGUNDY
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_chapter_slide(prs, chapter_num, title):
    """Add a chapter title slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BURGUNDY_DARK)

    # Chapter number
    chapter_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(12.333), Inches(1)
    )
    tf = chapter_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Chapitre {chapter_num}"
    p.font.size = Pt(24)
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(2.9),
        Inches(5.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Chapter title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2),
        Inches(12.333), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PARCHMENT
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    # Fleur-de-lis
    deco_box = slide.shapes.add_textbox(
        Inches(6), Inches(5.5),
        Inches(1.333), Inches(0.8)
    )
    tf = deco_box.text_frame
    p = tf.paragraphs[0]
    p.text = "\u269C"
    p.font.size = Pt(36)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, has_screenshot_placeholder=False):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, PARCHMENT)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BURGUNDY
    title_bar.line.fill.background()

    # Gold accent
    gold_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.15),
        Inches(13.333), Inches(0.08)
    )
    gold_accent.fill.solid()
    gold_accent.fill.fore_color.rgb = GOLD
    gold_accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"

    # Content area
    if has_screenshot_placeholder:
        content_width = 6
        content_left = 0.5
    else:
        content_width = 12
        content_left = 0.5

    content_box = slide.shapes.add_textbox(
        Inches(content_left), Inches(1.6),
        Inches(content_width), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"\u2726  {bullet}"  # Four-pointed star bullet
        p.font.size = Pt(20)
        p.font.color.rgb = INK
        p.font.name = "Georgia"
        p.space_after = Pt(12)
        p.level = 0

    if has_screenshot_placeholder:
        # Placeholder for screenshot
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7), Inches(1.6),
            Inches(5.8), Inches(5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = PARCHMENT_DARK
        placeholder.line.color.rgb = BURGUNDY
        placeholder.line.width = Pt(2)

        # Placeholder text
        ph_text = slide.shapes.add_textbox(
            Inches(7.5), Inches(3.8),
            Inches(4.8), Inches(0.8)
        )
        tf = ph_text.text_frame
        p = tf.paragraphs[0]
        p.text = "[Capture d'ecran]"
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = INK_LIGHT
        p.font.name = "Georgia"
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_feature_slide(prs, title, description, features):
    """Add a feature detail slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, PARCHMENT)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BURGUNDY
    title_bar.line.fill.background()

    # Gold accent
    gold_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.15),
        Inches(13.333), Inches(0.08)
    )
    gold_accent.fill.solid()
    gold_accent.fill.fore_color.rgb = GOLD
    gold_accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"

    # Description
    desc_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(12.333), Inches(1)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = INK_LIGHT
    p.font.name = "Georgia"

    # Features in cards
    card_width = 3.8
    card_height = 2.8
    start_x = 0.7
    start_y = 2.8
    gap = 0.3

    for i, (feat_title, feat_desc) in enumerate(features):
        col = i % 3
        row = i // 3

        x = start_x + col * (card_width + gap)
        y = start_y + row * (card_height + gap)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RgbColor(0xFF, 0xFF, 0xFF)
        card.line.color.rgb = BURGUNDY
        card.line.width = Pt(2)

        # Card left accent
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(0.1), Inches(card_height)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = BURGUNDY
        accent.line.fill.background()

        # Feature title
        feat_title_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.2),
            Inches(card_width - 0.4), Inches(0.6)
        )
        tf = feat_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = feat_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BURGUNDY
        p.font.name = "Georgia"

        # Feature description
        feat_desc_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.7),
            Inches(card_width - 0.4), Inches(card_height - 0.9)
        )
        tf = feat_desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feat_desc
        p.font.size = Pt(14)
        p.font.color.rgb = INK
        p.font.name = "Georgia"

    return slide

def add_screenshot_slide(prs, title, caption=""):
    """Add a slide for screenshots."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, PARCHMENT)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(1)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BURGUNDY
    title_bar.line.fill.background()

    # Gold accent
    gold_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.95),
        Inches(13.333), Inches(0.08)
    )
    gold_accent.fill.solid()
    gold_accent.fill.fore_color.rgb = GOLD
    gold_accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2),
        Inches(12.333), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"

    # Large placeholder for screenshot
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.4),
        Inches(10.333), Inches(5)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = PARCHMENT_DARK
    placeholder.line.color.rgb = BURGUNDY
    placeholder.line.width = Pt(3)

    # Placeholder text
    ph_text = slide.shapes.add_textbox(
        Inches(4), Inches(3.6),
        Inches(5.333), Inches(0.8)
    )
    tf = ph_text.text_frame
    p = tf.paragraphs[0]
    p.text = "[Inserer capture d'ecran ici]"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = INK_LIGHT
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    if caption:
        caption_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(6.6),
            Inches(10.333), Inches(0.6)
        )
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = INK_LIGHT
        p.font.name = "Georgia"
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_conclusion_slide(prs, title, points):
    """Add a conclusion slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BURGUNDY_DARK)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(12.333), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.font.name = "Georgia"
    p.alignment = PP_ALIGN.CENTER

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3), Inches(1.5),
        Inches(7.333), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Points
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(11.333), Inches(4.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"\u2726  {point}"
        p.font.size = Pt(22)
        p.font.color.rgb = PARCHMENT
        p.font.name = "Georgia"
        p.space_after = Pt(16)
        p.alignment = PP_ALIGN.CENTER

    # Fleur-de-lis
    deco_box = slide.shapes.add_textbox(
        Inches(6), Inches(6.5),
        Inches(1.333), Inches(0.8)
    )
    tf = deco_box.text_frame
    p = tf.paragraphs[0]
    p.text = "\u269C"
    p.font.size = Pt(36)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =====================
    # TITLE SLIDE
    # =====================
    add_title_slide(
        prs,
        "La Compagnie d'Excalibur",
        "Application de maintenance assistee par Intelligence Artificielle\n\nIAckathon 2024"
    )

    # =====================
    # CHAPTER 1: CHECKLIST
    # =====================
    add_chapter_slide(prs, "I", "Inspection des Destriers\nChecklist de Maintenance")

    # Checklist creation
    add_feature_slide(
        prs,
        "Creation de Checklist Dynamique",
        "Un formulaire interactif pour guider le technicien dans son inspection",
        [
            ("Formulaire Structure", "Checklist JSON dynamique avec sections et questions personnalisables"),
            ("Navigation Intuitive", "Progression par chapitres avec indicateur de progression visuel"),
            ("Numero de Serie", "Identification unique du vehicule (Destrier) avant inspection"),
            ("Pieces Jointes", "Capture photo directe depuis l'appareil pour documenter les defauts"),
            ("Commentaires", "Zone de notes optionnelle pour chaque question"),
            ("Sauvegarde Auto", "Enregistrement automatique en base de donnees locale"),
        ]
    )

    add_screenshot_slide(prs, "Formulaire de Checklist", "Interface de saisie avec navigation par sections")

    # AI Analysis
    add_feature_slide(
        prs,
        "Analyse IA des Defauts",
        "Le bouton 'Oracle' invoque l'intelligence artificielle pour analyser les photos",
        [
            ("Detection Automatique", "Analyse visuelle des photos pour identifier les defauts industriels"),
            ("Generation de Tags", "Creation automatique de mots-cles et descriptions des anomalies"),
            ("Modele Gemma Local", "IA executee localement sur l'appareil, sans connexion requise"),
            ("Anti-Hallucination", "Regles strictes pour eviter les faux positifs"),
            ("Bounding Box", "Localisation precise des defauts detectes sur l'image"),
            ("Confiance", "Score de confiance pour chaque detection"),
        ]
    )

    add_screenshot_slide(prs, "Bouton Oracle - Analyse IA", "Resultat de l'analyse avec tags generes automatiquement")

    # Completion & Tags
    add_content_slide(
        prs,
        "Completion et Edition des Tags",
        [
            "Resume automatique des reponses apportees",
            "Affichage des tags generes par l'IA (Sceaux Magiques)",
            "Possibilite de modifier ou supprimer les tags IA",
            "Ajout de tags manuels par le technicien",
            "Distinction visuelle entre tags IA et manuels",
            "Sauvegarde des tags modifies en base de donnees",
        ],
        has_screenshot_placeholder=True
    )

    # History
    add_feature_slide(
        prs,
        "Archives du Royaume",
        "Historique complet des checklists avec recherche par tags",
        [
            ("Liste Complete", "Toutes les checklists completees avec date et numero de serie"),
            ("Recherche Avancee", "Filtrage par tag IA, numero de serie ou titre"),
            ("Details Expansibles", "Vue detaillee des tags pour chaque checklist"),
            ("Statistiques", "Nombre de champs remplis et tags generes"),
            ("Navigation Directe", "Redirection automatique apres sauvegarde"),
            ("Refresh", "Actualisation de la liste par glissement"),
        ]
    )

    add_screenshot_slide(prs, "Archives - Historique des Checklists", "Liste des inspections avec tags IA")

    # =====================
    # CHAPTER 2: ASK PDF
    # =====================
    add_chapter_slide(prs, "II", "Consultation des Grimoires\nAsk my PDF")

    add_feature_slide(
        prs,
        "Ask my PDF - RAG Local",
        "Posez des questions sur vos documents techniques en langage naturel",
        [
            ("Import PDF", "Chargement de documents techniques depuis l'appareil"),
            ("Extraction Texte", "Analyse et indexation du contenu des documents"),
            ("Questions Naturelles", "Interface de chat pour poser des questions"),
            ("Reponses Contextuelles", "IA qui repond en citant les sources du document"),
            ("100% Hors-ligne", "Traitement local sans envoi de donnees"),
            ("Multi-Documents", "Possibilite d'interroger plusieurs PDF"),
        ]
    )

    add_screenshot_slide(prs, "Interface Ask PDF", "Chat avec le document technique")

    add_screenshot_slide(prs, "Reponse avec Sources", "L'IA cite les passages pertinents du PDF")

    # =====================
    # CONCLUSION
    # =====================
    add_conclusion_slide(
        prs,
        "Technologies Utilisees",
        [
            "Flutter / Dart - Application mobile cross-platform",
            "Gemma 3 Nano (4B) - Modele IA local quantifie int4",
            "SQLite / Drift - Base de donnees locale",
            "BLoC Pattern - Gestion d'etat reactive",
            "Material Design 3 - Interface utilisateur moderne",
            "Theme Kaamelott - Experience utilisateur immersive"
        ]
    )

    add_title_slide(
        prs,
        "Merci !",
        "La Compagnie d'Excalibur\n\n\u2694  Par le pouvoir de l'IA locale  \u2694"
    )

    # Save
    output_path = "presentation_iackathon.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
